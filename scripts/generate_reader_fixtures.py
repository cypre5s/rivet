"""生成 Phase 8 无版权、多格式且可重复的 Reader fixture。"""

from __future__ import annotations

import gzip
import importlib
import io
import json
import math
import os
import struct
import subprocess
import tarfile
import tempfile
import wave
import zipfile
from collections.abc import Callable
from datetime import UTC, datetime
from pathlib import Path
from typing import Protocol, cast

FIXTURE_ROOT = Path(__file__).parents[1] / "tests" / "fixtures" / "files"
FIXED_ZIP_TIME = (2026, 8, 28, 0, 0, 0)
FIXED_DOCUMENT_TIME = datetime(2026, 8, 28, tzinfo=UTC)
FREE_SECTOR = 0xFFFFFFFF
END_OF_CHAIN = 0xFFFFFFFE
FAT_SECTOR = 0xFFFFFFFD


class LegacyWorksheet(Protocol):
    """描述 fixture 生成所需的最小 xlwt Sheet API。"""

    def write(self, row: int, column: int, value: str) -> None:
        """写入一个文本单元格。"""
        ...


class LegacyWorkbook(Protocol):
    """描述 fixture 生成所需的最小 xlwt Workbook API。"""

    def add_sheet(self, name: str) -> LegacyWorksheet:
        """创建工作表。"""
        ...

    def save(self, path: str) -> None:
        """保存旧版 XLS。"""
        ...


def _zip_bytes(members: dict[str, bytes]) -> bytes:
    """以固定时间和顺序生成 ZIP 字节。"""
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name in sorted(members):
            info = zipfile.ZipInfo(name, date_time=FIXED_ZIP_TIME)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.create_system = 3
            info.external_attr = 0o100644 << 16
            archive.writestr(info, members[name])
    return output.getvalue()


def _normalize_zip(path: Path) -> None:
    """重写第三方库生成的 ZIP 容器以固定时间和成员顺序。"""
    with zipfile.ZipFile(path) as archive:
        members = {name: archive.read(name) for name in archive.namelist()}
    path.write_bytes(_zip_bytes(members))


def _pdf_bytes(text: str) -> bytes:
    """生成只含一个 Helvetica 文本流的最小有效 PDF。"""
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream = f"BT /F1 24 Tf 72 720 Td ({escaped}) Tj ET".encode("ascii")
    objects = (
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length "
        + str(len(stream)).encode("ascii")
        + b" >>\nstream\n"
        + stream
        + b"\nendstream",
    )
    document = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for object_number, body in enumerate(objects, start=1):
        offsets.append(len(document))
        document.extend(f"{object_number} 0 obj\n".encode("ascii"))
        document.extend(body)
        document.extend(b"\nendobj\n")
    xref_offset = len(document)
    document.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    document.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        document.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    document.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("ascii")
    )
    return bytes(document)


def _docx_bytes() -> bytes:
    """生成包含单段文字的最小 DOCX。"""
    members = {
        "[Content_Types].xml": b"""<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>""",
        "_rels/.rels": b"""<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>""",
        "word/document.xml": b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"><w:body><w:p><w:r><w:t>Rivet DOCX</w:t></w:r></w:p><w:sectPr/></w:body></w:document>""",
    }
    return _zip_bytes(members)


def _epub_bytes() -> bytes:
    """生成包含一个 XHTML 章节的 EPUB。"""
    members = {
        "META-INF/container.xml": b"""<?xml version="1.0"?>
<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container"><rootfiles><rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/></rootfiles></container>""",
        "OEBPS/chapter.xhtml": b"""<html xmlns="http://www.w3.org/1999/xhtml"><head><title>Rivet</title></head><body><p>Rivet EPUB</p></body></html>""",
        "OEBPS/content.opf": b"""<?xml version="1.0" encoding="UTF-8"?>
<package xmlns="http://www.idpf.org/2007/opf" version="3.0" unique-identifier="id"><metadata xmlns:dc="http://purl.org/dc/elements/1.1/"><dc:identifier id="id">rivet-fixture</dc:identifier><dc:title>Rivet EPUB</dc:title><dc:language>zh-CN</dc:language></metadata><manifest><item id="chapter" href="chapter.xhtml" media-type="application/xhtml+xml"/></manifest><spine><itemref idref="chapter"/></spine></package>""",
        "mimetype": b"application/epub+zip",
    }
    return _zip_bytes(members)


def _directory_entry(
    name: str,
    *,
    object_type: int,
    right_sibling: int = FREE_SECTOR,
    child: int = FREE_SECTOR,
    start_sector: int = END_OF_CHAIN,
    stream_size: int = 0,
) -> bytes:
    """生成一个 CFB v3 Directory Entry。"""
    encoded_name = (name + "\x00").encode("utf-16-le")
    if len(encoded_name) > 64:
        raise ValueError("CFB 名称过长")
    entry = bytearray(128)
    entry[: len(encoded_name)] = encoded_name
    struct.pack_into("<H", entry, 64, len(encoded_name))
    entry[66] = object_type
    entry[67] = 1
    struct.pack_into("<III", entry, 68, FREE_SECTOR, right_sibling, child)
    struct.pack_into("<I", entry, 116, start_sector)
    struct.pack_into("<Q", entry, 120, stream_size)
    return bytes(entry)


def _msg_bytes() -> bytes:
    """生成包含四个 Outlook 属性流的最小 CFB MSG。"""
    streams = (
        ("__substg1.0_0037001F", "Rivet MSG"),
        ("__substg1.0_0C1F001F", "fixture@example.invalid"),
        ("__substg1.0_0E04001F", "reader@example.invalid"),
        ("__substg1.0_1000001F", "Rivet MSG body"),
    )
    sector_count = 3 + len(streams) * 8
    fat = [FREE_SECTOR] * 128
    fat[0] = FAT_SECTOR
    fat[1] = 2
    fat[2] = END_OF_CHAIN
    stream_starts: list[int] = []
    next_sector = 3
    for _name, _content in streams:
        stream_starts.append(next_sector)
        for offset in range(8):
            sector = next_sector + offset
            fat[sector] = END_OF_CHAIN if offset == 7 else sector + 1
        next_sector += 8
    header = bytearray(512)
    header[:8] = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
    struct.pack_into("<HH", header, 24, 0x003E, 3)
    struct.pack_into("<HH", header, 28, 0xFFFE, 9)
    struct.pack_into("<H", header, 32, 6)
    struct.pack_into("<I", header, 44, 1)
    struct.pack_into("<I", header, 48, 1)
    struct.pack_into("<I", header, 56, 4096)
    struct.pack_into("<I", header, 60, END_OF_CHAIN)
    struct.pack_into("<I", header, 68, END_OF_CHAIN)
    struct.pack_into("<I", header, 76, 0)
    for index in range(1, 109):
        struct.pack_into("<I", header, 76 + index * 4, FREE_SECTOR)
    fat_sector = b"".join(struct.pack("<I", value) for value in fat)
    entries = [
        _directory_entry("Root Entry", object_type=5, child=1),
    ]
    for index, ((name, _content), start_sector) in enumerate(
        zip(streams, stream_starts, strict=True),
        start=1,
    ):
        right_sibling = index + 1 if index < len(streams) else FREE_SECTOR
        entries.append(
            _directory_entry(
                name,
                object_type=2,
                right_sibling=right_sibling,
                start_sector=start_sector,
                stream_size=4096,
            )
        )
    directory = b"".join(entries).ljust(1024, b"\x00")
    data_sectors = bytearray()
    for _name, content in streams:
        encoded = content.encode("utf-16-le")
        padding = (" " * ((4096 - len(encoded)) // 2)).encode("utf-16-le")
        data_sectors.extend((encoded + padding)[:4096].ljust(4096, b"\x00"))
    payload = bytes(header) + fat_sector + directory + bytes(data_sectors)
    if len(payload) != 512 + sector_count * 512:
        raise AssertionError("CFB 扇区数量不一致")
    return payload


def _write_office_files(root: Path) -> None:
    """使用格式官方 Python 原语生成 PPTX、XLSX 与 XLS。"""
    from openpyxl import Workbook
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    text_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(4), Inches(1))
    text_box.text = "Rivet PPTX"
    presentation.save(str(root / "sample.pptx"))
    _normalize_zip(root / "sample.pptx")

    workbook = Workbook()
    workbook.properties.created = FIXED_DOCUMENT_TIME
    workbook.properties.modified = FIXED_DOCUMENT_TIME
    sheet = workbook.active
    if sheet is None:
        raise AssertionError("XLSX 必须包含活动工作表")
    sheet.title = "Rivet"
    sheet.append(["name", "value"])
    sheet.append(["project", "Rivet XLSX"])
    workbook.save(root / "sample.xlsx")
    _normalize_zip(root / "sample.xlsx")

    xlwt = importlib.import_module("xlwt")
    workbook_factory = cast(Callable[[], LegacyWorkbook], xlwt.__dict__["Workbook"])
    legacy = workbook_factory()
    legacy_sheet = legacy.add_sheet("Rivet")
    legacy_sheet.write(0, 0, "name")
    legacy_sheet.write(0, 1, "value")
    legacy_sheet.write(1, 0, "project")
    legacy_sheet.write(1, 1, "Rivet XLS")
    legacy.save(str(root / "sample.xls"))


def _write_images(root: Path) -> None:
    """用 Pillow 生成六种小型真实图片。"""
    from PIL import Image

    image = Image.new("RGB", (64, 32), color=(32, 96, 160))
    formats = {
        "sample.png": "PNG",
        "sample.jpg": "JPEG",
        "sample.webp": "WEBP",
        "sample.gif": "GIF",
        "sample.bmp": "BMP",
        "sample.tiff": "TIFF",
    }
    for name, image_format in formats.items():
        image.save(root / name, format=image_format)


def _write_wav(root: Path) -> None:
    """生成 0.25 秒、8 kHz 单声道 PCM WAV。"""
    sample_rate = 8_000
    frames = bytearray()
    for index in range(sample_rate // 4):
        value = int(8_000 * math.sin(2 * math.pi * 440 * index / sample_rate))
        frames.extend(struct.pack("<h", value))
    with wave.open(str(root / "sample.wav"), "wb") as stream:
        stream.setnchannels(1)
        stream.setsampwidth(2)
        stream.setframerate(sample_rate)
        stream.writeframes(bytes(frames))


def _run_ffmpeg(arguments: list[str]) -> None:
    """使用 imageio 固定二进制执行一个无 shell 媒体生成命令。"""
    imageio_ffmpeg = importlib.import_module("imageio_ffmpeg")
    get_executable = cast(Callable[[], str], imageio_ffmpeg.__dict__["get_ffmpeg_exe"])

    subprocess.run(
        [
            get_executable(),
            "-hide_banner",
            "-loglevel",
            "error",
            "-y",
            *arguments,
        ],
        check=True,
        capture_output=True,
    )


def _write_media(root: Path) -> None:
    """生成五种音频与四种视频容器。"""
    exact_audio = ["-fflags", "+bitexact", "-flags:a", "+bitexact", "-threads", "1"]
    audio_outputs = {
        "sample.mp3": [*exact_audio, "-c:a", "libmp3lame", "-write_xing", "0"],
        "sample.m4a": [*exact_audio, "-c:a", "aac"],
        "sample.flac": [*exact_audio, "-c:a", "flac"],
        "sample.ogg": [*exact_audio, "-c:a", "libvorbis", "-serial_offset", "0"],
    }
    for name, codec in audio_outputs.items():
        _run_ffmpeg(
            [
                "-i",
                str(root / "sample.wav"),
                *codec,
                "-metadata",
                "creation_time=2026-08-28T00:00:00Z",
                str(root / name),
            ]
        )
    exact_video = ["-fflags", "+bitexact", "-flags:v", "+bitexact", "-threads", "1"]
    video_outputs = {
        "sample.mp4": [*exact_video, "-c:v", "libx264", "-pix_fmt", "yuv420p"],
        "sample.mov": [*exact_video, "-c:v", "libx264", "-pix_fmt", "yuv420p"],
        "sample.webm": [*exact_video, "-c:v", "libvpx-vp9"],
        "sample.mkv": [*exact_video, "-c:v", "libx264", "-pix_fmt", "yuv420p"],
    }
    for name, codec in video_outputs.items():
        _run_ffmpeg(
            [
                "-f",
                "lavfi",
                "-i",
                "color=c=blue:s=64x64:d=0.5:r=4",
                *codec,
                "-metadata",
                "creation_time=2026-08-28T00:00:00Z",
                str(root / name),
            ]
        )


def _write_archives(root: Path) -> None:
    """生成含同一文本成员的 ZIP、TAR、TGZ 与 7z。"""
    import py7zr

    member_content = b"Rivet archive\n"
    (root / "sample.zip").write_bytes(_zip_bytes({"inside.txt": member_content}))
    tar_stream = io.BytesIO()
    with tarfile.open(fileobj=tar_stream, mode="w") as archive:
        info = tarfile.TarInfo("inside.txt")
        info.size = len(member_content)
        info.mtime = 0
        info.mode = 0o644
        archive.addfile(info, io.BytesIO(member_content))
    tar_bytes = tar_stream.getvalue()
    (root / "sample.tar").write_bytes(tar_bytes)
    gzip_stream = io.BytesIO()
    with gzip.GzipFile(fileobj=gzip_stream, mode="wb", filename="", mtime=0) as stream:
        stream.write(tar_bytes)
    (root / "sample.tar.gz").write_bytes(gzip_stream.getvalue())
    with tempfile.TemporaryDirectory(prefix="rivet-fixture-7z-") as directory:
        member_path = Path(directory) / "inside.txt"
        member_path.write_bytes(member_content)
        os.utime(member_path, ns=(0, 0))
        with py7zr.SevenZipFile(root / "sample.7z", mode="w") as archive:
            archive.write(member_path, arcname="inside.txt")


def main() -> None:
    """写入固定清单中的全部 Reader fixture。"""
    FIXTURE_ROOT.mkdir(parents=True, exist_ok=True)
    text_files = {
        "sample.txt": "Rivet text\n",
        "sample.json": '{"project": "Rivet JSON"}\n',
        "sample.yaml": "project: Rivet YAML\n",
        "sample.toml": 'project = "Rivet TOML"\n',
        "sample.xml": "<root><project>Rivet XML</project></root>\n",
        "sample.csv": "name,value\nproject,Rivet CSV\n",
        "sample.tsv": "name\tvalue\nproject\tRivet TSV\n",
        "sample.html": "<html><body><h1>Rivet HTML</h1></body></html>\n",
        "sample.eml": (
            "From: fixture@example.invalid\n"
            "To: reader@example.invalid\n"
            "Subject: Rivet EML\n"
            "MIME-Version: 1.0\n"
            "Content-Type: text/plain; charset=utf-8\n\n"
            "Rivet EML body\n"
        ),
    }
    for name, content in text_files.items():
        (FIXTURE_ROOT / name).write_text(content, encoding="utf-8")
    notebook: dict[str, object] = {
        "nbformat": 4,
        "nbformat_minor": 5,
        "metadata": {},
        "cells": [
            {
                "cell_type": "markdown",
                "metadata": {},
                "source": ["Rivet notebook"],
            }
        ],
    }
    (FIXTURE_ROOT / "sample.ipynb").write_text(
        json.dumps(notebook, ensure_ascii=False, separators=(",", ":")) + "\n",
        encoding="utf-8",
    )
    (FIXTURE_ROOT / "sample.pdf").write_bytes(_pdf_bytes("Rivet PDF"))
    (FIXTURE_ROOT / "sample.docx").write_bytes(_docx_bytes())
    (FIXTURE_ROOT / "sample.epub").write_bytes(_epub_bytes())
    (FIXTURE_ROOT / "sample.msg").write_bytes(_msg_bytes())
    (FIXTURE_ROOT / "sample.bin").write_bytes(b"\x00\xffRIVET-BINARY-STRING\x00\x01")
    _write_office_files(FIXTURE_ROOT)
    _write_images(FIXTURE_ROOT)
    _write_wav(FIXTURE_ROOT)
    _write_media(FIXTURE_ROOT)
    _write_archives(FIXTURE_ROOT)


if __name__ == "__main__":
    main()
